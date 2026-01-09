# import os
# import uuid
# from typing import List, Dict, Optional, Tuple
# from datetime import datetime, timedelta
# import PyPDF2
# import pdfplumber
# import chromadb
# from sentence_transformers import SentenceTransformer
# import tiktoken
# from pathlib import Path
# from app.core.logging_config import logger

# class DocumentProcessingService:
#     """
#     Processes educational documents (PDFs, notes) for context-aware quiz generation.
#     Extracts content, chunks it, and stores in vector database.
#     """

    
#     def __init__(self):
#         # Initialize embedding model
#         self.embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
        
#         # Initialize ChromaDB (vector database)
#         self.chroma_client = chromadb.PersistentClient(path="./chroma_db")
#         self.collection = self.chroma_client.get_or_create_collection(
#             name="educational_documents",
#             metadata={"description": "TVET educational materials"}
#         )
        
#         # Token counter for chunking
#         self.tokenizer = tiktoken.get_encoding("cl100k_base")
        
#         # Chunking parameters
#         self.chunk_size = 500  # tokens per chunk
#         self.chunk_overlap = 50  # overlap between chunks
    
#     async def process_pdf(
#         self,
#         file_path: str,
#         metadata: Dict
#     ) -> Dict:
#         """
#         Extract text from PDF and store in vector database.
        
#         Args:
#             file_path: Path to PDF file
#             metadata: {
#                 "course": "Electrical Wiring",
#                 "topic": "Circuit Breakers",
#                 "week": 3,
#                 "upload_date": "2024-01-15",
#                 "instructor": "Mr. Kamau"
#             }
#         """
#         try:
#             logger.info(f"Processing PDF: {file_path}")
            
#             # Extract text from PDF
#             text_content = self._extract_pdf_text(file_path)
            
#             # Extract structured content (tables, sections)
#             structured_content = self._extract_structured_content(file_path)
            
#             # Chunk the content
#             chunks = self._chunk_text(text_content)
            
#             # Generate embeddings and store
#             doc_id = str(uuid.uuid4())
#             self._store_chunks(doc_id, chunks, metadata, structured_content)
            
#             logger.info(f"Successfully processed PDF: {doc_id}")
            
#             return {
#                 "document_id": doc_id,
#                 "file_name": Path(file_path).name,
#                 "total_chunks": len(chunks),
#                 "total_tokens": sum(len(self.tokenizer.encode(chunk)) for chunk in chunks),
#                 "metadata": metadata,
#                 "processed_at": datetime.now().isoformat(),
#                 "structured_content": {
#                     "sections": len(structured_content.get("sections", [])),
#                     "key_terms": len(structured_content.get("key_terms", [])),
#                     "has_diagrams": structured_content.get("has_diagrams", False)
#                 }
#             }
            
#         except Exception as e:
#             logger.error(f"PDF processing failed: {e}")
#             raise
    
#     def _extract_pdf_text(self, file_path: str) -> str:
#         """Extract text from PDF using multiple methods for robustness."""
#         text_content = []
        
#         # Method 1: PyPDF2 (fast, basic)
#         try:
#             with open(file_path, 'rb') as file:
#                 pdf_reader = PyPDF2.PdfReader(file)
#                 for page in pdf_reader.pages:
#                     text = page.extract_text()
#                     if text:
#                         text_content.append(text)
#         except Exception as e:
#             logger.warning(f"PyPDF2 extraction failed: {e}")
        
#         # Method 2: pdfplumber (better for complex layouts)
#         if not text_content:
#             try:
#                 with pdfplumber.open(file_path) as pdf:
#                     for page in pdf.pages:
#                         text = page.extract_text()
#                         if text:
#                             text_content.append(text)
#             except Exception as e:
#                 logger.error(f"pdfplumber extraction failed: {e}")
#                 raise
        
#         full_text = "\n\n".join(text_content)
#         return full_text
    
#     def _extract_structured_content(self, file_path: str) -> Dict:
#         """Extract structured information from PDF."""
#         structured = {
#             "sections": [],
#             "key_terms": [],
#             "tables": [],
#             "has_diagrams": False
#         }
        
#         try:
#             with pdfplumber.open(file_path) as pdf:
#                 for page_num, page in enumerate(pdf.pages):
#                     # Extract tables
#                     tables = page.extract_tables()
#                     if tables:
#                         structured["tables"].extend([
#                             {"page": page_num + 1, "data": table}
#                             for table in tables
#                         ])
                    
#                     # Detect images/diagrams
#                     if page.images:
#                         structured["has_diagrams"] = True
                    
#                     # Extract text and identify sections (headings)
#                     text = page.extract_text()
#                     if text:
#                         lines = text.split('\n')
#                         for line in lines:
#                             # Simple heuristic: lines in all caps or numbered are sections
#                             if line.isupper() and len(line.split()) <= 6:
#                                 structured["sections"].append({
#                                     "title": line.strip(),
#                                     "page": page_num + 1
#                                 })
                            
#                             # Extract potential key terms (words in bold, caps, or italics)
#                             # This is simplified - in production, use NER models
#                             words = line.split()
#                             for word in words:
#                                 if word.isupper() and len(word) > 3:
#                                     structured["key_terms"].append(word)
        
#         except Exception as e:
#             logger.warning(f"Structured extraction failed: {e}")
        
#         # Deduplicate key terms
#         structured["key_terms"] = list(set(structured["key_terms"]))[:20]
        
#         return structured

    






import os
import uuid
from typing import List, Dict, Optional
from datetime import datetime, timedelta
import PyPDF2
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
import chromadb
from sentence_transformers import SentenceTransformer
import tiktoken
from pathlib import Path
from app.core.logging_config import logger

class DocumentProcessingService:
    """
    Processes educational documents (PDFs, Word, PowerPoint, Text) 
    for context-aware quiz generation.
    """
    
    def __init__(self):
        self.embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
        self.chroma_client = chromadb.PersistentClient(path="./chroma_db")
        self.collection = self.chroma_client.get_or_create_collection(
            name="educational_documents",
            metadata={"description": "TVET educational materials"}
        )
        self.tokenizer = tiktoken.get_encoding("cl100k_base")
        self.chunk_size = 500
        self.chunk_overlap = 50
        
        # Supported file types
        self.supported_types = {
            '.pdf': self._extract_pdf_text,
            '.docx': self._extract_docx_text,
            '.doc': self._extract_docx_text,  # Try docx parser for .doc too
            '.pptx': self._extract_pptx_text,
            '.ppt': self._extract_pptx_text,
            '.txt': self._extract_txt_text,
            '.md': self._extract_txt_text,  # Markdown as text
        }
    
    async def process_document(
        self,
        file_path: str,
        metadata: Dict
    ) -> Dict:
        """
        Extract text from various document types and store in vector database.
        
        Supports: PDF, DOCX, DOC, PPTX, PPT, TXT, MD
        """
        try:
            file_ext = Path(file_path).suffix.lower()
            
            if file_ext not in self.supported_types:
                raise ValueError(f"Unsupported file type: {file_ext}")
            
            logger.info(f"Processing {file_ext} file: {file_path}")
            
            # Extract text using appropriate method
            extractor = self.supported_types[file_ext]
            text_content = extractor(file_path)
            
            if not text_content or len(text_content.strip()) < 100:
                raise ValueError("Document appears to be empty or too short")
            
            # Extract structured content (if applicable)
            structured_content = self._extract_structured_content_generic(
                file_path, file_ext, text_content
            )
            
            # Chunk the content
            chunks = self._chunk_text(text_content)
            
            # Generate embeddings and store
            doc_id = str(uuid.uuid4())
            metadata['file_type'] = file_ext
            metadata['original_filename'] = Path(file_path).name
            
            self._store_chunks(doc_id, chunks, metadata, structured_content)
            
            logger.info(f"Successfully processed {file_ext}: {doc_id}")
            
            return {
                "document_id": doc_id,
                "file_name": Path(file_path).name,
                "file_type": file_ext,
                "total_chunks": len(chunks),
                "total_tokens": sum(len(self.tokenizer.encode(chunk)) for chunk in chunks),
                "total_characters": len(text_content),
                "metadata": metadata,
                "processed_at": datetime.now().isoformat(),
                "structured_content": {
                    "sections": len(structured_content.get("sections", [])),
                    "key_terms": len(structured_content.get("key_terms", [])),
                    "has_images": structured_content.get("has_images", False)
                }
            }
            
        except Exception as e:
            logger.error(f"Document processing failed: {e}")
            raise
    
    def _extract_pdf_text(self, file_path: str) -> str:
        """Extract text from PDF."""
        text_content = []
        
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        text_content.append(text)
        except Exception as e:
            logger.warning(f"pdfplumber failed, trying PyPDF2: {e}")
            try:
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page in pdf_reader.pages:
                        text = page.extract_text()
                        if text:
                            text_content.append(text)
            except Exception as e2:
                logger.error(f"PDF extraction failed: {e2}")
                raise
        
        return "\n\n".join(text_content)
    
    def _extract_docx_text(self, file_path: str) -> str:
        """Extract text from Word documents (.docx, .doc)."""
        try:
            doc = DocxDocument(file_path)
            
            text_content = []
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_content.append(para.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells])
                    if row_text.strip():
                        text_content.append(row_text)
            
            full_text = "\n\n".join(text_content)
            logger.info(f"Extracted {len(full_text)} characters from Word document")
            return full_text
            
        except Exception as e:
            logger.error(f"Word document extraction failed: {e}")
            raise ValueError(f"Could not read Word document: {e}")
    
    def _extract_pptx_text(self, file_path: str) -> str:
        """Extract text from PowerPoint presentations (.pptx, .ppt)."""
        try:
            prs = Presentation(file_path)
            
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"--- Slide {slide_num} ---"]
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                    
                    # Extract text from tables in slides
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join([cell.text for cell in row.cells])
                            if row_text.strip():
                                slide_text.append(row_text)
                
                if len(slide_text) > 1:  # More than just the header
                    text_content.extend(slide_text)
            
            full_text = "\n\n".join(text_content)
            logger.info(f"Extracted text from {len(prs.slides)} slides")
            return full_text
            
        except Exception as e:
            logger.error(f"PowerPoint extraction failed: {e}")
            raise ValueError(f"Could not read PowerPoint: {e}")
    
    def _extract_txt_text(self, file_path: str) -> str:
        """Extract text from plain text files (.txt, .md)."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
            
            logger.info(f"Extracted {len(text)} characters from text file")
            return text
            
        except UnicodeDecodeError:
            # Try different encoding
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    text = file.read()
                return text
            except Exception as e:
                logger.error(f"Text file extraction failed: {e}")
                raise ValueError(f"Could not read text file: {e}")
    
    def _extract_structured_content_generic(
        self, 
        file_path: str, 
        file_ext: str, 
        text_content: str
    ) -> Dict:
        """Extract structured information from any document type."""
        structured = {
            "sections": [],
            "key_terms": [],
            "has_images": False,
            "tables": []
        }
        
        try:
            # Extract sections (lines that look like headers)
            lines = text_content.split('\n')
            for i, line in enumerate(lines):
                line = line.strip()
                
                # Heuristics for section headers
                if line and (
                    line.isupper() and len(line.split()) <= 8 or  # ALL CAPS
                    line.startswith('#') or  # Markdown header
                    (len(line) < 50 and i > 0 and lines[i-1].strip() == '') or  # Short line after blank
                    any(line.startswith(prefix) for prefix in ['Chapter', 'Section', 'Part', 'Unit'])
                ):
                    structured["sections"].append({
                        "title": line,
                        "line_number": i
                    })
                
                # Extract potential key terms (capitalized technical terms)
                words = line.split()
                for word in words:
                    # Technical terms are often capitalized or acronyms
                    if word.isupper() and len(word) > 2 and len(word) < 15:
                        structured["key_terms"].append(word)
            
            # Deduplicate and limit key terms
            structured["key_terms"] = list(set(structured["key_terms"]))[:30]
            
            # Check for images (for DOCX/PPTX)
            if file_ext in ['.docx', '.pptx']:
                structured["has_images"] = True  # Assume yes for these formats
            
            logger.info(f"Found {len(structured['sections'])} sections, {len(structured['key_terms'])} key terms")
            
        except Exception as e:
            logger.warning(f"Structured extraction failed: {e}")
        
        return structured
    
    
    
    def _chunk_text(self, text: str) -> List[str]:
        """Split text into overlapping chunks for better context."""
        # Tokenize
        tokens = self.tokenizer.encode(text)
        
        chunks = []
        start = 0
        
        while start < len(tokens):
            # Get chunk of tokens
            end = start + self.chunk_size
            chunk_tokens = tokens[start:end]
            
            # Decode back to text
            chunk_text = self.tokenizer.decode(chunk_tokens)
            chunks.append(chunk_text)
            
            # Move start forward (with overlap)
            start += self.chunk_size - self.chunk_overlap
        
        logger.info(f"Created {len(chunks)} chunks from text")
        return chunks
    
    def _store_chunks(
        self,
        doc_id: str,
        chunks: List[str],
        metadata: Dict,
        structured_content: Dict
    ):
        """Store text chunks with embeddings in ChromaDB."""
        
        # Generate embeddings for all chunks
        embeddings = self.embedding_model.encode(chunks).tolist()
        
        # Prepare data for storage
        ids = [f"{doc_id}_chunk_{i}" for i in range(len(chunks))]
        metadatas = [
            {
                **metadata,
                "chunk_index": i,
                "document_id": doc_id,
                "chunk_size": len(chunk),
                "has_key_terms": any(term in chunk for term in structured_content.get("key_terms", []))
            }
            for i, chunk in enumerate(chunks)
        ]
        
        # Store in ChromaDB
        self.collection.add(
            ids=ids,
            documents=chunks,
            embeddings=embeddings,
            metadatas=metadatas
        )
        
        logger.info(f"Stored {len(chunks)} chunks for document {doc_id}")
    
    def retrieve_relevant_content(
        self,
        query: str,
        filters: Optional[Dict] = None,
        top_k: int = 5
    ) -> List[Dict]:
        """
        Retrieve most relevant document chunks for a query.
        
        Args:
            query: Search query (e.g., "circuit breakers")
            filters: {"course": "Electrical Wiring", "week": 3}
            top_k: Number of chunks to retrieve
        """
        # Generate query embedding
        query_embedding = self.embedding_model.encode([query])[0].tolist()
        
        # Build where clause for filtering
        where = {}
        if filters:
            where = {k: v for k, v in filters.items() if v is not None}
        
        # Search in ChromaDB
        results = self.collection.query(
            query_embeddings=[query_embedding],
            n_results=top_k,
            where=where if where else None
        )
        
        # Format results
        relevant_chunks = []
        if results['documents']:
            for i in range(len(results['documents'][0])):
                relevant_chunks.append({
                    "content": results['documents'][0][i],
                    "metadata": results['metadatas'][0][i],
                    "distance": results['distances'][0][i],
                    "id": results['ids'][0][i]
                })
        
        logger.info(f"Retrieved {len(relevant_chunks)} relevant chunks for query: {query}")
        return relevant_chunks
    
    def get_documents_by_timeframe(
        self,
        course: str,
        timeframe: str = "daily"  # "daily", "weekly", "monthly"
    ) -> List[Dict]:
        """Get documents uploaded in specific timeframe."""
        
        # Calculate date range
        now = datetime.now()
        if timeframe == "daily":
            start_date = now - timedelta(days=1)
        elif timeframe == "weekly":
            start_date = now - timedelta(weeks=1)
        elif timeframe == "monthly":
            start_date = now - timedelta(days=30)
        else:
            start_date = now - timedelta(days=7)  # default to weekly
        
        # Query ChromaDB for documents in timeframe
        # Note: This is a simplified query - in production, you'd query by date metadata
        results = self.collection.get(
            where={"course": course}
        )
        
        # Group by document_id
        documents = {}
        if results['documents']:
            for i in range(len(results['documents'])):
                doc_id = results['metadatas'][i].get('document_id')
                if doc_id not in documents:
                    documents[doc_id] = {
                        "document_id": doc_id,
                        "metadata": results['metadatas'][i],
                        "chunks": []
                    }
                documents[doc_id]["chunks"].append(results['documents'][i])
        
        return list(documents.values())
    
    def get_document_summary(self, document_id: str) -> Dict:
        """Get summary of a processed document."""
        
        results = self.collection.get(
            where={"document_id": document_id}
        )
        
        if not results['documents']:
            return {"error": "Document not found"}
        
        # Aggregate information
        metadata = results['metadatas'][0]
        total_chunks = len(results['documents'])
        
        # Extract key information
        all_text = " ".join(results['documents'])
        
        return {
            "document_id": document_id,
            "metadata": metadata,
            "total_chunks": total_chunks,
            "total_characters": len(all_text),
            "preview": all_text[:500] + "...",
            "key_topics": metadata.get("topic", "Unknown")
        }